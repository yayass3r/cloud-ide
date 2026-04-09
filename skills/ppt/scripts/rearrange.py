#!/usr/bin/env python3
"""
Rearrange PowerPoint slides based on a sequence of indices.

Usage:
    python rearrange.py template.pptx output.pptx 0,34,34,50,52

Slides are 0-indexed. Indices can repeat to duplicate slides.
"""

import argparse
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


def copy_slide(src_prs: Presentation, dst_prs: Presentation, index: int, dst_layouts: dict) -> None:
    """Append a copy of slide[index] from src_prs into dst_prs."""
    src_slide = src_prs.slides[index]

    # Match layout by name across all masters; fall back to first available layout
    layout_name = src_slide.slide_layout.name
    dst_layout = dst_layouts.get(layout_name) or dst_prs.slide_layouts[0]

    new_slide = dst_prs.slides.add_slide(dst_layout)

    # Clear auto-added placeholder shapes
    for shape in list(new_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Collect image/media relationships from source
    image_rels = {
        rel_id: rel
        for rel_id, rel in src_slide.part.rels.items()
        if "image" in rel.reltype or "media" in rel.reltype
    }

    # Copy all shape elements
    for shape in src_slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        # Remap blip embed references to new slide's relationships
        r_embed = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
        for blip in new_el.xpath(".//a:blip[@r:embed]"):
            old_rId = blip.get(r_embed)
            if old_rId in image_rels:
                rel = image_rels[old_rId]
                new_rId = new_slide.part.rels.get_or_add(rel.reltype, rel._target)
                blip.set(r_embed, new_rId)

    # Copy slide-level background if defined.
    # p:bg lives inside p:cSld, not directly under p:sld.
    src_cSld = src_slide.element.find(qn("p:cSld"))
    dst_cSld = new_slide.element.find(qn("p:cSld"))
    if src_cSld is not None and dst_cSld is not None:
        src_bg = src_cSld.find(qn("p:bg"))
        if src_bg is not None:
            existing_bg = dst_cSld.find(qn("p:bg"))
            if existing_bg is not None:
                dst_cSld.remove(existing_bg)
            dst_cSld.insert(0, deepcopy(src_bg))


def rearrange_presentation(
    template_path: Path, output_path: Path, slide_sequence: list[int]
) -> None:
    src_prs = Presentation(template_path)
    total = len(src_prs.slides)

    for idx in slide_sequence:
        if idx < 0 or idx >= total:
            raise ValueError(f"Slide index {idx} out of range (0–{total - 1})")

    # Build a fresh presentation with the same dimensions
    dst_prs = Presentation(template_path)

    # Remove all existing slides from dst_prs
    sldIdLst = dst_prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))  # must use full namespace via qn(), not bare "r:id"
        if rId:
            dst_prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    # Search all slide masters for layout matching (templates may have multiple masters)
    all_layouts = {
        layout.name: layout
        for master in dst_prs.slide_masters
        for layout in master.slide_layouts
    }

    # Append slides in requested order (duplicates included)
    for idx in slide_sequence:
        copy_slide(src_prs, dst_prs, idx, all_layouts)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    dst_prs.save(output_path)
    print(f"Saved {len(slide_sequence)} slides → {output_path}")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Rearrange PowerPoint slides.",
        epilog="Example: python rearrange.py template.pptx output.pptx 0,34,34,50,52",
    )
    parser.add_argument("template", help="Path to template PPTX")
    parser.add_argument("output", help="Path for output PPTX")
    parser.add_argument("sequence", help="Comma-separated 0-based slide indices")
    args = parser.parse_args()

    template_path = Path(args.template)
    if not template_path.exists():
        print(f"Error: Template not found: {args.template}")
        sys.exit(1)

    try:
        slide_sequence = [int(x.strip()) for x in args.sequence.split(",")]
    except ValueError:
        print("Error: sequence must be comma-separated integers (e.g. 0,34,34,50,52)")
        sys.exit(1)

    try:
        rearrange_presentation(template_path, Path(args.output), slide_sequence)
    except ValueError as e:
        print(f"Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
